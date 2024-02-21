from pptx import Presentation
from pathlib import Path
from langchain.text_splitter import RecursiveCharacterTextSplitter


def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

    return "\n".join(text)


def read_all_pptx_files(root_folder_path):
    pptx_files = []

    for path in root_folder_path.rglob("*.pptx"):
        print("Load File: ", path)
        pptx_files.append(path)

    texts = []
    for pptx_file in pptx_files:
        text = extract_text_from_pptx(pptx_file)
        texts.append(text)

    return texts


def readtxt_file(path='Locator.txt'):
    with open(path, 'r') as f:
        text = f.read()
    return text


def text_split(state_of_the_union: str):
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=230, chunk_overlap=10)
    text_array = text_splitter.create_documents([state_of_the_union])
    str_array = [text.page_content for text in text_array]

    return str_array


def get_chunk():
    text = readtxt_file()
    text_list = text_split(text)
    return text_list


if __name__ == "__main__":
    # check chunk size and content
    text = readtxt_file()
    text_list = text_split(text)
    print(len(text_list))
